"""
Build an EDITABLE PowerPoint (.pptx) of the talk from the same content as slides_claude.md.

Real text boxes (editable) + placed figures from assets/ + speaker notes in each slide's
Notes field — so PowerPoint's Presenter View shows the notes on your laptop while the slides
show on the projector.

Run: .venv/bin/python presentation/make_pptx.py   ->  presentation/slides_claude.pptx
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "slides_claude.pptx")

NAVY = RGBColor(0x1F, 0x49, 0x7D)
BLUE = RGBColor(0x00, 0x70, 0xC0)
GREY = RGBColor(0x44, 0x44, 0x44)
DARK = RGBColor(0x22, 0x22, 0x22)

SW, SH = 13.333, 7.5  # 16:9 inches


def img_aspect(path):
    with Image.open(path) as im:
        w, h = im.size
    return w / h


def add_title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(SW - 1.1), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"


def add_bullets(slide, bullets, top=1.45, size=20):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(SW - 1.4), Inches(1.25))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run(); r.text = "•  " + b
        r.font.size = Pt(size); r.font.color.rgb = DARK; r.font.name = "Calibri"


def add_caption(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(SH - 0.6), Inches(SW - 1.2), Inches(0.45))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"


def place_images(slide, images, top, avail_h):
    n = len(images)
    gap = 0.3
    max_w_total = SW - 1.2
    per_w = (max_w_total - gap * (n - 1)) / n
    placed = []
    for path in images:
        asp = img_aspect(path)
        w = per_w; h = w / asp
        if h > avail_h:
            h = avail_h; w = h * asp
        placed.append((path, w, h))
    total_w = sum(w for _, w, _ in placed) + gap * (n - 1)
    x = (SW - total_w) / 2
    for path, w, h in placed:
        y = top + (avail_h - h) / 2
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        x += w + gap


def content_slide(title, bullets=None, images=None, caption=None, note=None):
    s = prs.slides.add_slide(BLANK)
    add_title(s, title)
    top = 1.45
    if bullets:
        add_bullets(s, bullets, top=top)
        top = 1.45 + 0.42 * len(bullets) + 0.35
    if images:
        bottom = (SH - 0.7) if caption else (SH - 0.35)
        place_images(s, [os.path.join(A, im) for im in images], top=top, avail_h=bottom - top)
    if caption:
        add_caption(s, caption)
    if note:
        s.notes_slide.notes_text_frame.text = note
    return s


def title_slide():
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(SW - 1.6), Inches(1.4))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Neural MCPF Allocation"
    r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
    sb = s.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(SW - 2.0), Inches(2.2))
    tf = sb.text_frame; tf.word_wrap = True
    for i, (txt, sz, col, it) in enumerate([
        ("Learning the Allocation Step of Multi-Agent Path Finding", 24, GREY, False),
        ("A transformer that replaces a combinatorial solver's allocation — "
         "near-optimal, any size, much faster at scale.", 18, BLUE, True),
        ("Ofek Yabo · Dayan Badalbaev — Multi-Agent Systems, 2026", 16, GREY, False),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(10)
        run = p.add_run(); run.text = txt
        run.font.size = Pt(sz); run.font.color.rgb = col; run.font.italic = it; run.font.name = "Calibri"
    s.notes_slide.notes_text_frame.text = (
        "One-sentence framing: we taught a neural network to make the expensive decision a classic "
        "MAPF solver makes — which agent visits which goals — and studied when that trade is worth it.")


prs = Presentation()
prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

title_slide()

content_slide(
    "Who goes where?",
    ["Many agents, many goals — minimize total work",
     "warehouse robots · delivery fleets · search & rescue"],
    images=["motivation.png"],
    note="The core decision in multi-agent coordination is the ALLOCATION: which agent takes which "
         "goals. A bad allocation wastes travel and creates collisions; real systems need it fast.")

content_slide(
    "The problem (formally)",
    ["mTSP inside MAPF: each goal visited once; agents take 0/1/many goals",
     "minimize total tour cost, then plan collision-free paths"],
    images=["bundle_vs_split.png"],
    note="Multi-agent TSP, not 1-to-1: an agent may take several goals if cheaper (bundle cost 7 beats "
         "split 14). This is why the model outputs a per-goal distribution over agents, not a permutation.")

content_slide(
    "The exact algorithm we build on",
    ["RobustMCPF — the optimal solver",
     "LKH-TSP allocation (k-best)  →  CBS collision-free planning  →  optimal paths",
     "exact / near-optimal, but the TSP allocation is combinatorial"],
    note="Ground truth = RobustMCPF (BasicMAPF). LKH finds the optimal allocation/tours; CBS makes "
         "collision-free paths. Quality is great; the allocation step is the badly-scaling part. "
         "Code: solver_wrapper.py → RobustMCPF/.")

content_slide(
    "The bottleneck — and our idea",
    ["LKH allocation grows with problem size; CBS is needed by both",
     "Idea: replace only the allocation with a neural net (near-instant after training)",
     "Expected: imitate the solver → big speedup, especially at scale",
     "Accept suboptimal solutions — the NN won't always find the optimum"],
    note="Swap the expensive combinatorial allocator for a learned one, keep everything else. "
         "Hypothesis: close to solver quality, large speedup as N,M grow.")

content_slide(
    "Where the NN drops in",
    images=["pipeline_solver_vs_nn.png"],
    caption="only the allocation box changes — CBS is shared (apples-to-apples)",
    note="Both pipelines share goal-ordering + CBS; only the allocator differs, so any cost/speed "
         "difference is attributable to the allocation. Code: evaluation/full_pipeline_eval.py.")

content_slide(
    "How we represent the problem",
    ["D (N×M): agent → goal   ·   G (M×M): goal → goal (normalized BFS)",
     "distances, not coordinates",
     "same distances as LKH-TSP — but two matrices (D, G), not one big cost table"],
    images=["D_G_matrices.png"],
    note="BFS distances respect walls and transfer across maps. G is the KEY feature — which goals are "
         "near each other (which to bundle); adding G lifted accuracy +14–18 points, more than any "
         "architecture/size/data change. Code: dataset_generation/distance.py.")

content_slide(
    "The network",
    ["row attention — each agent attends over its goals",
     "column attention — each goal attends over its agents",
     "per block: row × column attention + FFN, repeated ×L",
     "output: column softmax  ·  loss = CE + λ·MinSum"],
    images=["attention_sketch.png"],
    note="Each agent-goal pair is embedded; stacked row/column attention captures the combinatorial "
         "structure. Output is a per-goal distribution over agents. Trained with cross-entropy to "
         "imitate the solver + small MinSum term (λ=0.1). Code: model/network.py, model/losses.py.")

content_slide(
    "One model for ANY N, M",
    images=["anyNM_pipeline.png"],
    caption="no positional embeddings → the same weights run a 2×2 or a 50×100 problem",
    note="KEY slide. Every scalar D[i,j] → a d-dim token (N×M grid). G injected per goal by embedding "
         "each G[j,k] and SUMMING over k (sum-pool → any M). No positional embeddings (agents/goals "
         "unordered; attention handles variable length). Output: Linear(d→1) → N×M logits → column "
         "softmax → argmax per goal = the binary allocation matrix. Code: GoalAllocTransformerUniversal.")

# Build slides: advance with the arrow keys to walk through the pipeline one stage at a time.
for _k in range(1, 8):
    content_slide(
        "One model for ANY N, M — step by step",
        images=[f"anynm_step{_k}.png"],
        note="Advance with the arrow keys. Pre = data prep (instance → D,G). Net = shared-weight "
             "transformer: embed each scalar → ℝ^d token (N×M grid), row+column attention ×L, per-token "
             "logit ℝ^d→1, column softmax. Post = argmax decode. G is sum-pooled per goal and ADDED into "
             "the D tokens (context, not extra tokens) → output stays N×M. d and L are fixed; the token "
             "count N×M grows with the instance — no fixed-size vector.")

content_slide(
    "Method",
    ["data: solver labels D,G,Y   ·   train: imitate, mixed-size   ·   eval: offline + true CBS cost"],
    images=["method_flow.png"],
    note="Generate instances, run the solver, store D,G,Y (reject unreachable/over-budget). Train one "
         "universal model on many (N,M) shapes. Evaluate offline (accuracy) and — what matters — true "
         "execution cost (NN→ordering→CBS vs solver). Code: build_dataset.py, train.py, "
         "full_pipeline_eval.py, tests/.")

content_slide(
    "Small vs. big model",
    ["h128/L6 = 1.2M params   ·   h256/L8 = 6.3M",
     "preview: small wins in-distribution, big wins far out-of-distribution"],
    images=["params_bar.png"],
    note="Identical except width/depth. In-distribution the small model is better and faster (big "
         "overfits, Exp 12); at extreme extrapolation the big one generalizes better (Exp 15).")

content_slide(
    "Our models: architecture × training data (A / B / C)",
    ["A = original (small scale)  ·  B = trained on the real maps  ·  C = random-diverse grids",
     "suffix 1 = small net (h128, 1.2M)  ·  2 = big net (h256, 6.3M)"],
    images=["model_inventory.png"],
    note="Two axes: the net (h128 vs h256) and the training data. A = original small-scale model "
         "(random small grids). B = trained on the 4 real benchmark maps; C copies B's recipe but trains "
         "on random 32×32 grids (0–50% walls) — a controlled data swap. Train sizes: A ≈840k, B & C ≈720k.")

content_slide(
    "Results: near the solver, cheaper to run",
    ["cost ratio ~1.02–1.05 (≈1–6% above optimal)",
     "exact-cost match ~60–99% — both track M, nearly flat in N"],
    images=["heatmap_cost.png", "heatmap_match.png"],
    note="Across 28 small configs the NN is ~1–6% above optimal; large diverse model mean 1.020. Most "
         "disagreements are cost-equivalent ties (exact-cost match ≫ exact-allocation match). CBS "
         "dominates full-pipeline time, so single-instance speedup is modest; batched it's ~1000×.")

content_slide(
    "From random grids to real benchmark maps",
    images=["maps_panel.png"],
    caption="open → scattered → maze → rooms; only maze/room are truly structured",
    note="The 4 standard RobustMCPF benchmark maps. Only maze and room carry corridor/room structure — "
         "that distinction drives the next results.")

content_slide(
    "Results: generalization studies",
    images=["scissors_13b.png", "random_vs_fixed_bymap.png"],
    note="(a) small model zero-shot on real maps at small N/M = 1.019, but at large N/M fails (1.246 vs "
         "map-trained 1.048); scissors at M=30 — more agents hurt the old model (1.28→1.44), help the "
         "in-range one (1.10→1.05): generalizes across SHAPE, not SCALE. (b) random-grid training ties "
         "fixed-map everywhere except the maze (1.029 vs 1.127): random walls never make corridors. "
         "(c) at N≤50/M≤100 the big model extrapolates best (1.208 vs 1.244), speedup 5.9–9.4×.")

content_slide(
    "Pushing further: zero-shot XL extrapolation (Exp 15)",
    ["B models run zero-shot at N∈{20,35,50} × M∈{50,75,100} (no retraining)",
     "the bigger model (h256) wins out-of-distribution (1.208 vs 1.244), and inference is "
     "5.9–9.4× faster than the solver"],
    images=["xl_heatmaps_speedup.png"],
    caption="green = closer to optimal · h256 better across the grid (esp. N=50) · M=75 is the hard column",
    note="Zero-shot XL toward the paper's N≤50/M≤100. Verdict flips vs in-distribution (Exp 12): the "
         "capacity that overfit in range is the better generalizer at scale — h256 1.208 vs h128 1.244, "
         "worst-case gap roughly halved (224→124). M=75 hardest for both; maze easiest (≈1.15–1.18). "
         "Speedup 5.9–9.4× as LKH's TSP grows. Numbers: RESULTS.md Exp 15.")

content_slide(
    "When to use the NN vs. the exact solver",
    images=["decision_table.png"],
    note="NN replaces only allocation, so it wins where allocation/TSP cost dominates and many instances "
         "are solved: big problems (5.9–9.4× at ~20% overhead), batched inference, unstructured maps. "
         "The exact solver still wins on tiny/exactness-critical/unseen-structured instances.")

content_slide(
    "Demo: same instance, solver vs NN",
    images=["demo_run_compare.gif"],
    caption="different allocation, identical cost (62) — the NN found its own equally-optimal plan",
    note="Real run of model A on the room map. Both reach cost 62, but the NN sends the two right-side "
         "goals to a different agent than the solver — CBS executes it collision-free at the same cost. "
         "A vivid example of the cost-equivalent ties behind the high exact-match rates. (GIF animates "
         "in PowerPoint slideshow; static side-by-side is assets/demo_nn_vs_solver.png.)")

content_slide(
    "Conclusion",
    ["Problem: optimal allocation (mTSP in MAPF) — combinatorial step scales badly",
     "Idea: a small universal transformer imitates it from distances — near-optimal",
     "Generalizes: any N,M; across scale and (with the right data) geometry",
     "Payoff: 5.9–9.4× at scale, within ~20% of optimal — use where speed/scale matter"],
    note="Recap the 4 messages.")

content_slide(
    "Thank you — questions?",
    ["Acknowledgment: the RobustMCPF solver (LKH + CBS) is the work of Yehonatan Kidushim "
     "(github.com/yehonatan280198)",
     "repo: Neural_MCPF_Allocation · RobustMCPF (LKH + CBS)"],
    note="Keep the backup slides ready for likely questions. Credit RobustMCPF (LKH+CBS) to Yehonatan Kidushim.")

# ---- extra / backup slides ----
content_slide(
    "(Extra) Loss",
    ["L = L_CE + λ·L_MinSum   (λ = 0.1)",
     "CE — per-goal cross-entropy vs the solver (imitation)",
     "MinSum = Σ P·D — mild geometric regularizer"],
    note="CE drives imitation; MinSum prevents degenerate uniform predictions early. Code: model/losses.py.")

content_slide(
    "(Extra) Overfitting & small-vs-big",
    images=["convergence.png"],
    note="Val loss bottoms ~epoch 22 then rises; best-val-loss checkpointing = effective early stopping. "
         "Explains why the big model overfits in-distribution but generalizes better far-OOD.")

content_slide(
    "(Extra) Metrics",
    ["Offline — allocation agreement (per-goal / full-assignment accuracy)",
     "Execution cost — true CBS path-length ratio (NN / solver)",
     "Exact match — identical integer cost (≫ exact-allocation match → ties)"],
    note="Code: evaluation/evaluate.py (offline), full_pipeline_eval.py (execution cost).")

n_slides = len(prs.slides._sldIdLst)
prs.save(OUT)
print("wrote", OUT, f"({n_slides} slides)")
